import pandas as pd 
import os 
from pptx import Presentation
from docx import Document
import pymupdf4llm
import fitz
import pymupdf
from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter, TextSplitter
import logging
from bs4 import BeautifulSoup
from langchain_text_splitters import HTMLHeaderTextSplitter
from langchain_text_splitters import RecursiveCharacterTextSplitter
import re
import nltk
import pytesseract
from PIL import Image
from docx import Document as DocxDocument
from langchain_community.document_loaders import (
    UnstructuredExcelLoader,
    TextLoader,
    CSVLoader,
    PyMuPDFLoader
)
chunk_size = 4000
chunk_overlap =150
nltk.download('punkt_tab')
nltk.download('averaged_perceptron_tagger_eng')

# Configure Tesseract path
# pytesseract.pytesseract.tesseract_cmd = 'C:/Program Files/Tesseract-OCR/tesseract.exe'

def chunk(data):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=150)
    splits = text_splitter.split_documents(data)
    return splits

def chunks(data):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=150)
    splits = text_splitter.create_documents(data)
    return splits

def chunk_text(text, chunk_size=4000):
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]


# Function to extract text from an image
def extract_text_from_image(image_path):
    print("path_img",image_path)
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text


def get_file_byte_string(blob_service_client, container_name, blob_name):

    logging.info("Initiating Container Client")
    container_client = blob_service_client.get_container_client(container_name)

    logging.info("Initiating Blob Client")
    blob_client = container_client.get_blob_client(blob_name)

    BLOB_FILE_EXIST = blob_client.exists()
    logging.info(f"Blob file exists: {BLOB_FILE_EXIST}")

    if BLOB_FILE_EXIST:
        logging.info("Downloading Blob")
        download_stream = blob_client.download_blob()
        FILE_CONTENT = download_stream.readall()
        logging.info("Blob Downloaded")
        return  FILE_CONTENT
    return None
def get_chunk_data(path):
    # pdf
    if path.endswith(".pdf"):

        pdf_doc = fitz.open(path)  # Changed to directly use fitz from PyMuPDF
        markdown_data = pymupdf4llm.to_markdown(pdf_doc)
        
        headers_to_split_on = [
            ("#", "Header 1"),
            ("##", "Header 2"),
            ("###", "Header 3"),
        ]
        chunk_size =4000
        chunk_overlap =150
        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on, strip_headers=False)
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size, chunk_overlap=chunk_overlap
        )

        markdown_splits = markdown_splitter.split_text(markdown_data)
        text_splits = text_splitter.split_documents(markdown_splits)
        chunks = [ts.page_content for ts in text_splits]
        return chunks

#  Ppptx
    elif path.endswith(".pptx"):
        presentation = Presentation(path)
        slide_data = []
        for slide in presentation.slides:
            # Access slide content, shapes, and text
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text)
            slide_text = "\n".join(slide_text)
            if slide_text:
                # print("slide_text", slide_text)
                slide_data.append(slide_text)
        return slide_data

# .docx
    elif path.endswith((".doc", ".docx")):

        print("path",path)
        doc = Document(path)
        doc_text = []
        for para in doc.paragraphs:
            doc_text.append(para.text)

        doc_text = "\n".join(doc_text)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=150)
        text_splits = text_splitter.split_text(doc_text)
        return text_splits
    
    # Excel Files
    elif path.endswith(".xlsx"):
        xl_data = []
        loader = UnstructuredExcelLoader(path)
        data = loader.load()
        data = chunk(data)
        for x in data:
            xl_data.append(x.page_content)
        return xl_data
    
    # CSV Files
    elif path.endswith(".csv"):
        csv_data = []
        loader = CSVLoader(path)
        data = loader.load()
        datas = chunk(data)
        for x in datas:
            csv_data.append(x.page_content)
        return csv_data
    
    # TXT Files
    elif path.endswith(".txt"):
        txt_data =[]
        loader = TextLoader(path)
        data = loader.load()
        data = chunk(data)
        for x in data:
            txt_data.append(x.page_content)
        return txt_data

    # Image Files
    elif path.endswith((".png", ".jpg", ".jpeg")):
        text = extract_text_from_image(path)   # n
        ap_text=""
        for x in text:
            ap_text+=x
        text = chunk_text(text,chunk_size = 4000)
        print("path",path)
        print('text",text')
        return text
# html
    elif path.endswith(".html"):

        def convert_table_contents_to_json(table):
            # Extract table headers
            headers = [header.text for header in table.find_all('th')]

            # Extract table rows
            rows = table.find_all('tr')

            # Initialize an empty list to store table data
            table_data = []

            # Loop through rows and extract data
            for row in rows:
                cells = row.find_all('td')
                if cells:
                    row_data = {headers[i]: cell.text for i, cell in enumerate(cells)}
                    table_data.append(row_data)

            json_data = json.dumps(table_data, indent=1)
            return json_data

        def extract_and_replace_tables(bs_object):
            """Extracts tables from bs_object and replaces them with comma seperated values.
            Returns a beasuiful soup object."""

            tables = bs_object.find_all('table')
            for table in tables:
                # csv_rows = convert_table_contents_to_csv(table)
                csv_rows = convert_table_contents_to_json(table)
                # print(csv_rows)
                table.replace_with(BeautifulSoup(csv_rows, 'html.parser'))

            return bs_object

        def extract_and_replace_a_tags(bs_object):
            a_tags = bs_object.find_all('a')
            # replace a tags with their text contents ans url i.e content (url)
            for a_tag in a_tags:
                if a_tag.text and a_tag.get('href'):
                    url = a_tag.get('href')
                    # if not url.startswith('http'):
                    #     base_url = main_url.split('//')[0] + '//' + main_url.split('//')[1].split('/')[0]
                    #     url = base_url + url
                    a_tag.replace_with(a_tag.text + ' (' + url + ')')
            return bs_object

        headers_to_split_on = [
            ("h1", "H 1"),
            ("h2", "H 2"),
            ("h3", "H 3"),
        ]

        html_splitter = HTMLHeaderTextSplitter(headers_to_split_on)
        # def get_html_chunks(html_bytes):
        soup = BeautifulSoup(path, 'html.parser')
        html_string = extract_and_replace_a_tags(soup)
        html_string = extract_and_replace_tables(soup)
        html_string = str(html_string)
        html_header_splits = html_splitter.split_text(html_string)
        html_chunks = []
        for html_split in html_header_splits:
            chunk_size = 4000
            chunk_overlap = 300
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size, chunk_overlap=chunk_overlap
            )
            splits = text_splitter.split_text(html_split.page_content)
            for split in splits:
                # attach metadat to each split
                html_chunks.append(str(html_split.metadata) +  " " + split)
        return html_chunks
    # Excel Files

    elif path.endswith(".xlsx"):
        xl_data = []
        loader = UnstructuredExcelLoader(path)
        data = loader.load()
        data = chunk(data)
        for x in data:
            xl_data.append(x.page_content)
        return xl_data
    
def get_direct_chunks(data):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=150)
    text_splits = text_splitter.split_text(data)
    return text_splits
    
def get_clean_id(text):
    # replace anything with _ other than alphanumeric characters
    text = re.sub(r"[^a-zA-Z0-9]+", '_', text)
    return text
